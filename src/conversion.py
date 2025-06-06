# -*- coding: utf-8 -*-

##########################################################################
# OpenLP - Open Source Lyrics Projection                                 #
# ---------------------------------------------------------------------- #
# Copyright (c) 2008 OpenLP Developers                                   #
# ---------------------------------------------------------------------- #
# This program is free software: you can redistribute it and/or modify   #
# it under the terms of the GNU General Public License as published by   #
# the Free Software Foundation, either version 3 of the License, or      #
# (at your option) any later version.                                    #
#                                                                        #
# This program is distributed in the hope that it will be useful,        #
# but WITHOUT ANY WARRANTY; without even the implied warranty of         #
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the          #
# GNU General Public License for more details.                           #
#                                                                        #
# You should have received a copy of the GNU General Public License      #
# along with this program.  If not, see <https://www.gnu.org/licenses/>. #
##########################################################################
"""
The :mod:`~openlp.plugins.mcp.conversion` module contains the ConversionWorker
class for background PowerPoint to PDF conversion.
"""

import logging
from pathlib import Path

from PySide6 import QtCore

log = logging.getLogger(__name__)


class ConversionWorker(QtCore.QObject):
    """Background worker for PowerPoint conversion to avoid GUI lockup."""
    conversion_completed = QtCore.Signal(object, str)  # (converted_path, title)
    conversion_failed = QtCore.Signal(str)  # error_message
    
    def __init__(self):
        super().__init__()
        
    def convert_powerpoint(self, ppt_path, title):
        """Convert PowerPoint in background thread."""
        try:
            # Move the conversion logic here
            pdf_path = self._convert_ppt_to_pdf(ppt_path)
            if pdf_path and pdf_path.exists():
                self.conversion_completed.emit(pdf_path, title)
            else:
                self.conversion_failed.emit(f"Failed to convert PowerPoint file to PDF")
        except Exception as e:
            self.conversion_failed.emit(f"Error during conversion: {str(e)}")
    
    def _convert_ppt_to_pdf(self, ppt_path):
        """Convert PowerPoint file to PDF using LibreOffice (preferred) or python-pptx fallback."""
        # First try LibreOffice (much better quality)
        libreoffice_result = self._convert_with_libreoffice(ppt_path)
        if libreoffice_result and libreoffice_result.exists():
            log.info(f"Successfully converted {ppt_path.name} using LibreOffice")
            return libreoffice_result
        
        # Fallback to python-pptx method
        log.info(f"LibreOffice not available, falling back to python-pptx for {ppt_path.name}")
        return self._convert_with_python_pptx(ppt_path)
    
    def _convert_with_libreoffice(self, ppt_path):
        """Convert PowerPoint file to PDF using LibreOffice."""
        try:
            import subprocess
            import time
            
            # Create temporary PDF file
            timestamp = int(time.time())
            pdf_name = f"{ppt_path.stem}_libreoffice_{timestamp}.pdf"
            pdf_path = ppt_path.parent / pdf_name
            
            # Try different LibreOffice command variations
            libreoffice_commands = [
                "soffice",
                "libreoffice", 
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/bin/soffice",
                "/usr/local/bin/soffice"
            ]
            
            for cmd in libreoffice_commands:
                try:
                    # Test if this command exists (quick check)
                    test_result = subprocess.run([cmd, "--version"], 
                                               capture_output=True, 
                                               text=True, 
                                               timeout=5)
                    if test_result.returncode == 0:
                        log.info(f"Found LibreOffice at: {cmd}, starting conversion...")
                        
                        # Perform the conversion with longer timeout
                        convert_cmd = [
                            cmd,
                            "--headless",
                            "--convert-to", "pdf",
                            "--outdir", str(ppt_path.parent),
                            str(ppt_path)
                        ]
                        
                        result = subprocess.run(convert_cmd, 
                                              capture_output=True, 
                                              text=True, 
                                              timeout=90)  # Increased to 90 seconds
                        
                        if result.returncode == 0:
                            # LibreOffice creates filename.pdf, rename to our unique name
                            default_pdf = ppt_path.with_suffix('.pdf')
                            if default_pdf.exists():
                                try:
                                    default_pdf.rename(pdf_path)
                                    log.info(f"LibreOffice conversion completed successfully")
                                    return pdf_path
                                except Exception:
                                    # Return the default PDF if rename fails
                                    log.info(f"LibreOffice conversion completed (using default name)")
                                    return default_pdf
                            elif pdf_path.exists():
                                return pdf_path
                        else:
                            log.debug(f"LibreOffice conversion failed: {result.stderr}")
                            
                except subprocess.TimeoutExpired:
                    log.warning(f"LibreOffice conversion timed out after 90 seconds")
                    return None
                except Exception as e:
                    log.debug(f"LibreOffice command {cmd} failed: {e}")
                    continue
            
            return None
                
        except Exception as e:
            log.debug(f"LibreOffice conversion error: {e}")
            return None
    
    def _convert_with_python_pptx(self, ppt_path):
        """Convert PowerPoint file to PDF using python-pptx and reportlab (fallback)."""
        try:
            # Try using python-pptx to extract content and create PDF
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import time
            
            # Create temporary PDF file
            timestamp = int(time.time())
            pdf_name = f"{ppt_path.stem}_python_{timestamp}.pdf"
            pdf_path = ppt_path.parent / pdf_name
            
            # Load PowerPoint presentation
            prs = Presentation(str(ppt_path))
            
            # Create PDF
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            width, height = letter
            
            for i, slide in enumerate(prs.slides):
                # Add slide number
                c.drawString(50, height - 50, f"Slide {i + 1}")
                
                y_position = height - 100
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Simple text extraction
                        text_lines = shape.text.strip().split('\n')
                        for line in text_lines:
                            if line.strip() and y_position > 50:
                                c.drawString(50, y_position, line[:80])  # Limit line length
                                y_position -= 20
                
                c.showPage()
            
            c.save()
            
            if pdf_path.exists():
                log.info(f"Successfully converted {ppt_path.name} to PDF using python-pptx")
                return pdf_path
            else:
                log.error("PDF conversion failed - file not created")
                return None
                
        except ImportError as e:
            log.warning("PowerPoint conversion requires python-pptx and reportlab: pip install python-pptx reportlab")
            return None
        except Exception as e:
            log.error(f"Error converting PowerPoint to PDF: {e}")
            return None 